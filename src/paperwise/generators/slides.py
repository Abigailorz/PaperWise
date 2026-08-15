"""PPT 生成 v3 — LLM 内容 + 确定性渲染。

内容层：SlideContentBuilder 用 LLM 生成结构化 slide JSON。
渲染层：SlideDeckRenderer 把 slide JSON 渲染成美观的 .pptx。
回退层：build_fallback_slides 在无 LLM / LLM 失败时用确定性提取兜底。
"""

import json
import re
from collections import deque
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

try:
    from PIL import Image
except Exception:  # pragma: no cover
    Image = None


# ==================== 设计令牌 ====================

COLORS = {
    "primary": RGBColor(0x10, 0x22, 0x36),
    "primary2": RGBColor(0x1C, 0x3D, 0x63),
    "accent": RGBColor(0x3A, 0x7B, 0xC2),
    "accent_soft": RGBColor(0xE8, 0xF1, 0xF9),
    "orange": RGBColor(0xC4, 0x82, 0x28),
    "dark": RGBColor(0x1A, 0x24, 0x32),
    "gray": RGBColor(0x5E, 0x6A, 0x77),
    "light": RGBColor(0xF2, 0xF5, 0xF8),
    "paper": RGBColor(0xFB, 0xFA, 0xF6),
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "green": RGBColor(0x2E, 0x7D, 0x5B),
    "green_soft": RGBColor(0xE3, 0xF4, 0xEA),
    "red": RGBColor(0xB2, 0x3B, 0x47),
    "red_soft": RGBColor(0xF9, 0xEA, 0xEC),
}

SERIES_COLORS = [
    COLORS["primary"], COLORS["accent"], COLORS["orange"],
    COLORS["green"], COLORS["red"], COLORS["primary2"],
]

FONT_TITLE = "Microsoft YaHei"
FONT_BODY = "Microsoft YaHei"

SLIDE_W = 13.333
SLIDE_H = 7.5


# ==================== 内容层 ====================

class SlideContentBuilder:
    """用 LLM 生成结构化 slide JSON。"""

    def __init__(self, llm_client):
        self.llm = llm_client

    async def build(self, title: str, paper_text: str,
                    report_sections: dict) -> dict:
        """生成 slide deck JSON。"""
        # 拼接可用上下文
        ctx_parts = [f"论文标题：{title}"]
        if paper_text:
            ctx_parts.append(f"论文正文（截断）：\n{paper_text[:12000]}")
        for name, content in (report_sections or {}).items():
            if content:
                ctx_parts.append(f"[{name}]\n{str(content)[:5000]}")
        context = "\n\n".join(ctx_parts)
        system = self._skill_system(paper_text, title)

        prompt = f"""你是资深学术 PPT 设计师。请把下面这篇论文整理成一份 10–15 页的演示文稿大纲。

严格输出一个 JSON 对象（不要输出任何 JSON 之外的文字、注释或 markdown 代码块）：
{{
  "title": "论文标题",
  "authors": "作者",
  "venue": "会议/期刊",
  "year": "年份",
  "slides": [
    {{"type": "title"}},
    {{"type": "summary", "headline": "一句话总结", "bullets": ["要点1", "要点2"],
      "stats": [{{"value": "199×", "label": "比 LERF 快"}}]}},
    {{"type": "section", "heading": "背景与动机"}},
    {{"type": "diagram", "heading": "方法总览", "layout": "TD",
      "mermaid": "graph TD\\nA[消息提取] --> B[正交 Procrustes 对齐] --> C[前缀注入]"}},
    {{"type": "content", "heading": "核心方法", "subheading": "正交 Procrustes 对齐",
      "bullets": ["要点（≤18字）"], "key_point": "关键数字或结论"}},
    {{"type": "chart", "heading": "主要结果", "kind": "bar",
      "categories": ["ARC-C", "MedQA", "GSM8K"],
      "series": [{{"name": "基线", "data": [93.2, 68.0, 87.6]}},
                {{"name": "本方法", "data": [93.7, 70.3, 89.8]}}]}},
    {{"type": "table", "heading": "消融实验", "headers": ["方法", "指标"], "rows": [["A", "0.9"]]}},
    {{"type": "flow", "heading": "系统流程", "steps": ["解析论文", "抽取要点", "生成解读"]}},
    {{"type": "grid", "heading": "核心贡献", "columns": 2,
      "items": [{{"title": "免训练", "text": "无需微调"}}]}},
    {{"type": "two_by_two", "heading": "优势与局限",
      "top_left": {{"title": "优势", "text": "..."}},
      "top_right": {{"title": "机会", "text": "..."}},
      "bottom_left": {{"title": "劣势", "text": "..."}},
      "bottom_right": {{"title": "威胁", "text": "..."}}}},
    {{"type": "quote", "text": "核心贡献一句话", "attribution": "作者"}},
    {{"type": "thanks"}}
  ]
}}

要求：
1. 以论文的论证为主线（为什么重要 → 瓶颈/缺口 → 做了什么 → 关键证据 → 为什么可信 → 边界与开放问题），不要照抄论文章节顺序；AI/方法类论文默认按 problem → solution 弧线组织；
2. 每页只讲一个核心信息，优先用图示表达，不要整页都是 bullet；结果页不要纯文字，有数据就用 chart、有流程就用 diagram；
3. 每个 bullet 中文 ≤ 18 字，英文 ≤ 12 词，务必精简，不要整句粘贴；控制每页文字量，避免溢出；
4. 必须把论文里的关键数字放进 stats / key_point / table / chart（如速度提升倍数、准确率、指标），图表数据必须来自论文真实数值，不得编造；
5. 方法页/架构页优先用全宽 diagram 画流程（mermaid 只用 graph TD 或 graph LR + A[标签] --> B[标签]，节点标签 ≤ 10 字）；实验结果页优先用 chart 画柱状图（kind=bar/hbar/line，data 为纯数字，不带 %）；贡献页用 grid；步骤页用 flow；对比页用 two_column 或 two_by_two；
6. 版式要有节奏变化，不要所有内容页都用同一种卡片模板，按证据类型选择版式；
7. slides 总数 10–15 页，其中 diagram/chart/flow/grid/two_by_two 至少出现 3 种。

论文内容：
{context}

请只输出 JSON。"""

        for attempt in range(2):
            try:
                resp = await self.llm.chat(
                    messages=[
                        {"role": "system", "content": system},
                        {"role": "user", "content": prompt},
                    ],
                    temperature=0.3, max_tokens=6000,
                )
                deck = self._parse_json(resp.content)
                if deck and deck.get("slides"):
                    deck.setdefault("title", title)
                    return deck
            except Exception:
                if attempt == 1:
                    break
        # 回退到确定性生成
        return build_fallback_slides({
            "title": title,
            "sections": report_sections,
            "overview": paper_text,
        })

    @staticmethod
    def _skill_system(paper_text: str, title: str) -> str:
        """加载 nature-paper2ppt skill 作为系统提示；失败时回退到内置提示。"""
        try:
            from paperwise.generators.pptx_skill import load_pptx_skill_prompt
            loaded = load_pptx_skill_prompt(paper_text, title)
            if loaded:
                return loaded
        except Exception:
            pass
        return "你是资深学术 PPT 设计师。请把论文整理成结构清晰、图示优先、图表数据真实的学术演示文稿。"

    @staticmethod
    def _parse_json(text: str) -> Optional[dict]:
        if not text:
            return None
        text = text.strip()
        # 去掉可能的 ```json ... ``` 包裹
        m = re.search(r"```(?:json)?\s*(\{.*\})\s*```", text, re.DOTALL)
        if m:
            text = m.group(1)
        try:
            return json.loads(text)
        except json.JSONDecodeError:
            m = re.search(r"\{.*\}", text, re.DOTALL)
            if m:
                try:
                    return json.loads(m.group(0))
                except json.JSONDecodeError:
                    return None
        return None


# ==================== 确定性回退 ====================

def _extract_points(text: str, max_points: int = 6) -> list[str]:
    if not text:
        return []
    text = re.sub(r"```.*?```", " ", text, flags=re.DOTALL)
    text = re.sub(r"!\[[^\]]*\]\([^)]*\)", " ", text)
    text = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", text)
    text = re.sub(r"#{1,6}\s*", "", text)
    lines = [l.strip() for l in text.split("\n") if l.strip()]
    raw = []
    for line in lines:
        if len(line) > 200:
            raw.extend(re.split(r"(?<=[。.!?;；])\s*", line))
        else:
            raw.append(line)
    points, seen = [], set()
    for r in raw:
        r = re.sub(r"^\s*(?:[-*•·>]|\d+[.、)])\s*", "", r).strip()
        r = re.sub(r"\s+", " ", r)
        if not (6 <= len(r) <= 160):
            continue
        key = r[:40]
        if key in seen:
            continue
        seen.add(key)
        points.append(r[:160])
        if len(points) >= max_points:
            break
    return points


def build_fallback_slides(paper_data: dict) -> dict:
    """无 LLM 时的确定性回退。"""
    title = paper_data.get("title") or "论文解读"
    sections = paper_data.get("sections") or {}
    overview = str(paper_data.get("overview") or sections.get("overview") or "")
    slides = [{"type": "title"}]
    if overview:
        slides.append({
            "type": "summary",
            "headline": "论文概览",
            "bullets": _extract_points(overview, max_points=5),
        })
    order = [
        ("motivation", "背景与动机"),
        ("methodology", "核心方法"),
        ("experiments", "实验与结果"),
        ("related_work", "相关工作"),
        ("conclusion", "总结与展望"),
    ]
    for key, heading in order:
        content = str(sections.get(key) or "").strip()
        if not content:
            continue
        slides.append({"type": "section", "heading": heading})
        slides.append({
            "type": "content",
            "heading": heading,
            "bullets": _extract_points(content, max_points=6),
        })
    crit = str(sections.get("critical_analysis") or "").strip()
    if crit:
        slides.append({
            "type": "two_column",
            "heading": "优势 vs 局限",
            "left": {"title": "优势", "items": []},
            "right": {"title": "局限", "items": _extract_points(crit, max_points=6)},
        })
    slides.append({"type": "thanks"})
    return {
        "title": title,
        "authors": paper_data.get("authors") or "",
        "venue": paper_data.get("venue") or "",
        "year": paper_data.get("year") or "",
        "slides": slides,
    }


# ==================== 渲染层 ====================

class SlideDeckRenderer:
    """把 slide JSON 渲染成 .pptx。"""

    def __init__(self, base_dir: Optional[Path] = None):
        self.prs = Presentation()
        self.prs.slide_width = Inches(SLIDE_W)
        self.prs.slide_height = Inches(SLIDE_H)
        self._base = Path(base_dir) if base_dir else None
        self._slide_no = 0

    def render(self, deck: dict, output_path: str) -> str:
        slides = deck.get("slides") or []
        for slide in slides:
            stype = slide.get("type", "content")
            handler = getattr(self, f"_render_{stype}", None) or self._render_content
            handler(slide, deck)
        out = Path(output_path)
        out.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(out))
        return str(out)

    # ---------- 各类型 ----------

    def _render_title(self, s, deck):
        slide = self._blank()
        self._bg(slide, COLORS["primary"])
        self._rect(slide, 0, 0, SLIDE_W, 0.14, COLORS["orange"], MSO_SHAPE.RECTANGLE)
        self._text(slide, 0.9, 2.0, 11.5, 2.6, deck.get("title", ""),
                   size=36, color=COLORS["white"], bold=True, align=PP_ALIGN.LEFT)
        sub = []
        if deck.get("authors"):
            sub.append(str(deck["authors"]))
        if deck.get("venue") or deck.get("year"):
            sub.append(" · ".join(x for x in [str(deck.get("venue") or ""), str(deck.get("year") or "")] if x))
        if sub:
            self._text(slide, 0.9, 4.7, 11.5, 1.2, "\n".join(sub),
                       size=18, color=RGBColor(0xBF, 0xD1, 0xF5), align=PP_ALIGN.LEFT)
        self._footer(slide, "PaperWise · AI 论文解读")

    def _render_section(self, s, deck):
        slide = self._blank()
        self._bg(slide, COLORS["primary"])
        self._rect(slide, 0.9, 3.25, 1.6, 0.09, COLORS["orange"], MSO_SHAPE.RECTANGLE)
        self._text(slide, 0.9, 2.7, 11.5, 1.2, str(s.get("heading") or ""),
                   size=42, color=COLORS["white"], bold=True, align=PP_ALIGN.LEFT)
        if s.get("subheading"):
            self._text(slide, 0.9, 3.7, 11.5, 0.8, str(s["subheading"]),
                       size=20, color=RGBColor(0xBF, 0xD1, 0xF5), align=PP_ALIGN.LEFT)
        self._footer(slide, str(deck.get("title") or ""))

    def _render_summary(self, s, deck):
        slide = self._content_header(str(s.get("headline") or "一句话总结"))
        bullets = s.get("bullets") or []
        self._bullets(slide, bullets, 0.9, 1.8, 7.0, 5.2, size=18)
        stats = s.get("stats") or []
        self._stat_cards(slide, stats, 8.1, 1.9, 4.4, 5.0)

    def _render_content(self, s, deck):
        heading = str(s.get("heading") or "内容")
        slide = self._content_header(heading)
        if s.get("subheading"):
            self._text(slide, 0.9, 1.52, 11.5, 0.5, str(s["subheading"]),
                       size=16, color=COLORS["gray"], align=PP_ALIGN.LEFT)
        top = 2.15 if s.get("subheading") else 1.85
        bullets = s.get("bullets") or []
        self._bullets(slide, bullets, 0.9, top, 9.0, 5.0, size=18)
        if s.get("key_point"):
            self._key_point_badge(slide, str(s["key_point"]), 10.1, top, 2.5, 4.8)

    def _render_two_column(self, s, deck):
        slide = self._content_header(str(s.get("heading") or "对比"))
        left = s.get("left") or {}
        right = s.get("right") or {}
        self._card_column(slide, 0.9, 1.9, 5.5, 5.0,
                          str(left.get("title") or "优势"), left.get("items") or [],
                          COLORS["green"], COLORS["green_soft"])
        self._card_column(slide, 7.0, 1.9, 5.5, 5.0,
                          str(right.get("title") or "局限"), right.get("items") or [],
                          COLORS["red"], COLORS["red_soft"])

    def _render_table(self, s, deck):
        slide = self._content_header(str(s.get("heading") or "结果"))
        headers = [str(h)[:40] for h in (s.get("headers") or [])]
        rows = s.get("rows") or []
        if not headers:
            return
        n_cols = min(len(headers), 6)
        n_rows = min(len(rows) + 1, 8)
        frame = slide.shapes.add_table(n_rows, n_cols, Inches(0.8), Inches(1.8),
                                       Inches(11.7), Inches(0.5 * n_rows))
        tbl = frame.table
        for c, h in enumerate(headers[:n_cols]):
            cell = tbl.cell(0, c)
            cell.text = h
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLORS["primary"]
            self._style_cell(cell, bold=True, color=COLORS["white"], size=14)
        for r, row in enumerate(rows[:n_rows - 1], start=1):
            for c in range(n_cols):
                val = row[c] if c < len(row) else ""
                cell = tbl.cell(r, c)
                cell.text = str(val)[:60]
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLORS["light"] if r % 2 else COLORS["white"]
                self._style_cell(cell, bold=False, color=COLORS["dark"], size=13)

    def _render_figure(self, s, deck):
        slide = self._content_header(str(s.get("heading") or "图示"))
        path = s.get("image") or s.get("path") or ""
        p = Path(path)
        if not p.is_absolute() and self._base:
            p = self._base / p
        if not str(path) or not p.exists():
            self._text(slide, 0.9, 2.5, 11.5, 1.5, "（未找到配图）",
                       size=18, color=COLORS["gray"], align=PP_ALIGN.CENTER)
            return
        self._picture_fit(slide, p)
        if s.get("caption"):
            self._text(slide, 1.0, 6.7, 11.3, 0.5, str(s["caption"]),
                       size=12, color=COLORS["gray"], align=PP_ALIGN.CENTER, italic=True)

    def _render_quote(self, s, deck):
        slide = self._blank()
        self._bg(slide, COLORS["light"])
        self._text(slide, 1.5, 2.3, 10.3, 2.0, f"“{s.get('text', '')}”",
                   size=30, color=COLORS["primary"], bold=True, align=PP_ALIGN.CENTER)
        if s.get("attribution"):
            self._text(slide, 1.5, 4.6, 10.3, 0.6, f"— {s['attribution']}",
                       size=16, color=COLORS["gray"], align=PP_ALIGN.CENTER)

    def _render_thanks(self, s, deck):
        slide = self._blank()
        self._bg(slide, COLORS["primary"])
        self._text(slide, 1.0, 2.8, 11.3, 1.6, "谢谢观看", size=48,
                   color=COLORS["white"], bold=True, align=PP_ALIGN.CENTER)
        self._text(slide, 1.0, 4.4, 11.3, 0.8, "Q & A", size=24,
                   color=RGBColor(0xBF, 0xD1, 0xF5), align=PP_ALIGN.CENTER)
        self._footer(slide, "由 PaperWise 生成")

    # ---------- 视觉型版式 ----------

    def _render_diagram(self, s, deck):
        slide = self._content_header(str(s.get("heading") or "示意图"))
        nodes, order, edges = self._parse_mermaid(s.get("mermaid") or s.get("diagram") or "")

        # 也支持直接传入 nodes / edges（结构化形式）
        for n in (s.get("nodes") or []):
            nid = str(n.get("id") or n.get("name") or "")
            if not nid:
                continue
            if nid not in nodes:
                nodes[nid] = str(n.get("label") or nid)
                order.append(nid)
        for e in (s.get("edges") or []):
            if not e or len(e) < 2:
                continue
            a, b = str(e[0]), str(e[1])
            lab = str(e[2]) if len(e) > 2 and e[2] else None
            nodes.setdefault(a, a.replace("_", " "))
            nodes.setdefault(b, b.replace("_", " "))
            if a not in order:
                order.append(a)
            if b not in order:
                order.append(b)
            edges.append((a, b, lab))

        if not nodes:
            self._text(slide, 1, 2.5, 11, 1, "（缺少流程图数据）",
                       size=18, color=COLORS["gray"], align=PP_ALIGN.CENTER)
            return

        layout = str(s.get("layout") or "TD").upper()
        groups = self._layout_diagram(nodes, order, edges, layout)
        positions = self._diagram_positions(groups, layout)
        dims = {nid: (self._est_w(label), 0.62) for nid, label in nodes.items()}

        # 先画箭头（在节点下层），再画节点覆盖
        for a, b, lab in edges:
            if a not in positions or b not in positions:
                continue
            x1, y1 = positions[a]
            x2, y2 = positions[b]
            w1, h1 = dims.get(a, (1.2, 0.62))
            w2, h2 = dims.get(b, (1.2, 0.62))
            dx, dy = x2 - x1, y2 - y1
            dist = (dx * dx + dy * dy) ** 0.5
            if dist < 1e-6:
                continue
            ux, uy = dx / dist, dy / dist
            sx, sy = self._clip_rect(x1, y1, w1 / 2, h1 / 2, ux, uy)
            ex, ey = self._clip_rect(x2, y2, w2 / 2, h2 / 2, -ux, -uy)
            self._arrow(slide, sx, sy, ex, ey, COLORS["gray"], 2.0)
            if lab:
                mx, my = (sx + ex) / 2, (sy + ey) / 2
                self._text(slide, mx - 1.5, my - 0.22, 3.0, 0.4, str(lab),
                           size=10, color=COLORS["gray"], align=PP_ALIGN.CENTER)

        for nid, label in nodes.items():
            if nid not in positions:
                continue
            cx, cy = positions[nid]
            w, h = dims[nid]
            self._node(slide, cx, cy, w, h, label)

    def _render_chart(self, s, deck):
        slide = self._content_header(str(s.get("heading") or "数据"))
        kind = str(s.get("kind") or "bar").lower()
        categories = [str(c) for c in (s.get("categories") or [])]
        series = s.get("series") or []
        if not categories or not series:
            self._text(slide, 1, 2.5, 11, 1, "（缺少图表数据）",
                       size=18, color=COLORS["gray"], align=PP_ALIGN.CENTER)
            return
        self._legend(slide, series, 1.2, 1.32, 11.0)
        if kind in ("hbar", "barh", "horizontal"):
            self._hbar_chart(slide, categories, series)
        elif kind == "line":
            self._line_chart(slide, categories, series)
        else:
            self._bar_chart(slide, categories, series)

    def _render_flow(self, s, deck):
        slide = self._content_header(str(s.get("heading") or "流程"))
        steps = [str(x) for x in (s.get("steps") or [])][:6]
        if not steps:
            self._text(slide, 1, 2.5, 11, 1, "（缺少流程步骤）",
                       size=18, color=COLORS["gray"], align=PP_ALIGN.CENTER)
            return
        n = len(steps)
        gap = 0.85
        step_w = min(2.4, (11.0 - gap * (n - 1)) / n)
        used = step_w * n + gap * (n - 1)
        x0 = (SLIDE_W - used) / 2
        top, h = 2.55, 1.5
        for i, st in enumerate(steps):
            x = x0 + i * (step_w + gap)
            card = self._rect(slide, x, top, step_w, h,
                              COLORS["primary"] if i % 2 == 0 else COLORS["primary2"])
            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.08)
            tf.margin_right = Inches(0.08)
            tf.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            p.text = str(i + 1)
            self._run_font(p.runs[0] if p.runs else p.add_run(), 18, COLORS["white"], True)
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.CENTER
            p2.text = st
            self._run_font(p2.runs[0] if p2.runs else p2.add_run(), 13, COLORS["white"], False)
            if i < n - 1:
                self._arrow(slide, x + step_w + 0.04, top + h / 2,
                            x + step_w + gap - 0.04, top + h / 2, COLORS["orange"], 2.2)

    def _render_grid(self, s, deck):
        slide = self._content_header(str(s.get("heading") or "要点"))
        items = s.get("items") or []
        if not items:
            self._text(slide, 1, 2.5, 11, 1, "（缺少要点数据）",
                       size=18, color=COLORS["gray"], align=PP_ALIGN.CENTER)
            return
        cols = max(1, min(int(s.get("columns") or 2), 3))
        items = items[:cols * 3]
        rows = (len(items) + cols - 1) // cols
        x0, y0, gap = 0.9, 1.9, 0.3
        total_w = 11.5
        card_w = (total_w - gap * (cols - 1)) / cols
        card_h = min(2.2, 4.7 / rows)
        for i, it in enumerate(items):
            r, c = divmod(i, cols)
            x = x0 + c * (card_w + gap)
            y = y0 + r * (card_h + gap)
            self._info_card(slide, x, y, card_w, card_h,
                            str(it.get("title") or ""),
                            str(it.get("text") or it.get("desc") or ""),
                            COLORS["primary2"], COLORS["light"])

    def _render_two_by_two(self, s, deck):
        slide = self._content_header(str(s.get("heading") or "二维对比"))
        x0, y0, cw, ch, gap = 0.9, 1.85, 5.65, 2.3, 0.3
        slots = {
            "top_left": (x0, y0),
            "top_right": (x0 + cw + gap, y0),
            "bottom_left": (x0, y0 + ch + gap),
            "bottom_right": (x0 + cw + gap, y0 + ch + gap),
        }
        for key, (x, y) in slots.items():
            q = s.get(key) or {}
            soft = COLORS["accent_soft"] if "left" in key else COLORS["light"]
            color = COLORS["primary"] if "left" in key else COLORS["accent"]
            self._info_card(slide, x, y, cw, ch,
                            str(q.get("title") or ""),
                            str(q.get("text") or ""),
                            color, soft)

    # ---------- 辅助 ----------

    def _blank(self):
        self._slide_no += 1
        return self.prs.slides.add_slide(self.prs.slide_layouts[6])

    def _bg(self, slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _rect(self, slide, x, y, w, h, color, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
        r = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
        r.fill.solid()
        r.fill.fore_color.rgb = color
        r.line.fill.background()
        r.shadow.inherit = False
        return r

    def _ellipse(self, slide, x, y, w, h, color):
        e = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
        e.fill.solid()
        e.fill.fore_color.rgb = color
        e.line.fill.background()
        e.shadow.inherit = False
        return e

    def _arrow(self, slide, x1, y1, x2, y2, color=COLORS["gray"], width=2.0):
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        conn.line.color.rgb = color
        conn.line.width = Pt(width)
        conn.shadow.inherit = False
        ln = conn.line._get_or_add_ln()
        tail = ln.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
        ln.append(tail)
        return conn

    def _line(self, slide, x1, y1, x2, y2, color=COLORS["accent"], width=2.0):
        conn = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        conn.line.color.rgb = color
        conn.line.width = Pt(width)
        conn.shadow.inherit = False
        return conn

    def _node(self, slide, cx, cy, w, h, text, fill=COLORS["primary"], color=COLORS["white"], size=14):
        r = self._rect(slide, cx - w / 2, cy - h / 2, w, h, fill)
        tf = r.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.08)
        tf.margin_right = Inches(0.08)
        tf.margin_top = Inches(0.02)
        tf.margin_bottom = Inches(0.02)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = text
        p.alignment = PP_ALIGN.CENTER
        self._run_font(p.runs[0] if p.runs else p.add_run(), size, color, True)
        return r

    def _info_card(self, slide, x, y, w, h, title, text, color, soft):
        card = self._rect(slide, x, y, w, h, soft)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.18)
        tf.margin_right = Inches(0.12)
        tf.margin_top = Inches(0.12)
        p = tf.paragraphs[0]
        p.text = title
        self._run_font(p.runs[0] if p.runs else p.add_run(), 17, color, True)
        if text:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(4)
            p2.text = text
            self._run_font(p2.runs[0] if p2.runs else p2.add_run(), 13, COLORS["dark"], False)
        return card

    def _legend(self, slide, series, x, y, w):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.36))
        tf = box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        for i, sr in enumerate(series[:6]):
            run = p.add_run()
            run.text = "● " + str(sr.get("name") or f"系列{i + 1}") + "   "
            self._run_font(run, 12, SERIES_COLORS[i % len(SERIES_COLORS)], False)

    def _bar_chart(self, slide, categories, series):
        plot_x, plot_y, plot_w, plot_h = 1.3, 2.05, 10.4, 4.15
        base_y = plot_y + plot_h
        max_v = self._series_max(series)
        self._line(slide, plot_x, base_y, plot_x + plot_w, base_y, COLORS["gray"], 1.2)
        self._line(slide, plot_x, plot_y, plot_x, base_y, COLORS["gray"], 1.2)
        n_cats = len(categories)
        n_ser = len(series)
        group_w = plot_w / n_cats
        bar_w = min(0.9, group_w * 0.62 / n_ser)
        show_labels = n_cats * n_ser <= 10
        for si, sr in enumerate(series[:6]):
            color = SERIES_COLORS[si % len(SERIES_COLORS)]
            data = sr.get("data") or []
            for ci in range(n_cats):
                v = self._num(data[ci]) if ci < len(data) else None
                if v is None:
                    continue
                gx = plot_x + group_w * ci + group_w / 2
                bx = gx - (n_ser * bar_w) / 2 + si * bar_w
                h = plot_h * max(v, 0) / max_v
                self._rect(slide, bx, base_y - h, bar_w, h, color, MSO_SHAPE.RECTANGLE)
                if show_labels:
                    self._text(slide, bx - 0.3, base_y - h - 0.32, bar_w + 0.6, 0.3,
                               f"{v:.1f}", size=9, color=COLORS["dark"], align=PP_ALIGN.CENTER)
        for ci, cat in enumerate(categories):
            gx = plot_x + group_w * ci + group_w / 2
            self._text(slide, gx - group_w / 2, base_y + 0.08, group_w, 0.35, cat,
                       size=10, color=COLORS["gray"], align=PP_ALIGN.CENTER)

    def _hbar_chart(self, slide, categories, series):
        plot_x, plot_y, plot_w, plot_h = 3.2, 1.95, 8.9, 4.35
        max_v = self._series_max(series)
        self._line(slide, plot_x, plot_y, plot_x, plot_y + plot_h, COLORS["gray"], 1.2)
        n_cats = len(categories)
        n_ser = len(series)
        row_h = plot_h / n_cats
        bar_h = min(0.35, row_h * 0.55 / n_ser)
        for ci, cat in enumerate(categories):
            cy = plot_y + row_h * ci + row_h / 2
            self._text(slide, 0.35, cy - 0.15, 2.75, 0.35, cat,
                       size=11, color=COLORS["dark"], align=PP_ALIGN.RIGHT)
            for si, sr in enumerate(series[:6]):
                data = sr.get("data") or []
                v = self._num(data[ci]) if ci < len(data) else None
                if v is None:
                    continue
                color = SERIES_COLORS[si % len(SERIES_COLORS)]
                by = cy - (n_ser * bar_h) / 2 + si * bar_h
                w = plot_w * max(v, 0) / max_v
                self._rect(slide, plot_x, by, w, bar_h, color, MSO_SHAPE.RECTANGLE)
                self._text(slide, plot_x + w + 0.06, by - 0.02, 1.0, bar_h,
                           f"{v:.1f}", size=9, color=COLORS["dark"])

    def _line_chart(self, slide, categories, series):
        plot_x, plot_y, plot_w, plot_h = 1.3, 2.05, 10.4, 4.15
        base_y = plot_y + plot_h
        max_v = self._series_max(series)
        self._line(slide, plot_x, base_y, plot_x + plot_w, base_y, COLORS["gray"], 1.2)
        self._line(slide, plot_x, plot_y, plot_x, base_y, COLORS["gray"], 1.2)
        n_cats = len(categories)
        if n_cats > 1:
            xs = [plot_x + plot_w * i / (n_cats - 1) for i in range(n_cats)]
        else:
            xs = [plot_x + plot_w / 2]
        for si, sr in enumerate(series[:6]):
            color = SERIES_COLORS[si % len(SERIES_COLORS)]
            data = sr.get("data") or []
            pts = []
            for ci in range(n_cats):
                v = self._num(data[ci]) if ci < len(data) else None
                if v is None:
                    continue
                pts.append((xs[ci], base_y - plot_h * max(v, 0) / max_v))
            for i in range(len(pts) - 1):
                self._line(slide, pts[i][0], pts[i][1], pts[i + 1][0], pts[i + 1][1], color, 2.2)
            for px, py in pts:
                self._ellipse(slide, px - 0.05, py - 0.05, 0.1, 0.1, color)
        for ci, cat in enumerate(categories):
            self._text(slide, xs[ci] - 0.6, base_y + 0.08, 1.2, 0.35, cat,
                       size=10, color=COLORS["gray"], align=PP_ALIGN.CENTER)

    @staticmethod
    def _series_max(series):
        vals = []
        for sr in series:
            for v in (sr.get("data") or []):
                n = SlideDeckRenderer._num(v)
                if n is not None:
                    vals.append(n)
        return max(vals) if vals else 1

    @staticmethod
    def _num(v):
        if isinstance(v, (int, float)):
            return float(v)
        s = str(v).strip().replace("%", "").replace("×", "").replace("x", "").replace("X", "")
        s = re.sub(r"[^\d.\-]", "", s)
        try:
            return float(s)
        except Exception:
            return None

    @staticmethod
    def _est_w(label, size=14):
        w = 0.0
        for ch in str(label):
            w += 0.21 if ord(ch) > 0x2E80 else 0.115
        return max(1.1, min(3.0, w * size / 14 + 0.55))

    @staticmethod
    def _clip_rect(cx, cy, hw, hh, dx, dy):
        tx = hw / abs(dx) if abs(dx) > 1e-6 else float("inf")
        ty = hh / abs(dy) if abs(dy) > 1e-6 else float("inf")
        t = min(tx, ty)
        return cx + dx * t, cy + dy * t

    def _parse_mermaid(self, spec):
        text = str(spec)
        nodes, order, edges = {}, [], []
        node_pat = re.compile(
            r"([A-Za-z0-9_\-]+)\s*(?:\[\[(.+?)\]\]|\[(.+?)\]|\(\((.+?)\)\)|\((.+?)\))")
        for m in node_pat.finditer(text):
            nid = m.group(1)
            label = next(g for g in m.groups()[1:] if g is not None)
            if nid not in nodes:
                nodes[nid] = label.strip()
                order.append(nid)
        arrow_re = re.compile(r"\s*(?:-->|---|-\.->)\s*")
        for raw in re.split(r"[\n;]", text):
            line = raw.strip()
            if not line:
                continue
            if re.match(r"^(graph|flowchart)\s+(TD|TB|LR|RL|BT)\b", line, re.I):
                continue
            low = line.lower()
            if low.startswith(("subgraph", "end", "style ", "classdef ", "class ", "linkstyle ")):
                continue
            stripped = node_pat.sub(lambda m: m.group(1), line)
            parts = [p.strip() for p in arrow_re.split(stripped) if p.strip()]
            for i in range(len(parts) - 1):
                src = parts[i]
                rest = parts[i + 1]
                label = None
                if rest.startswith("|"):
                    lm = re.match(r"\|([^|]*)\|\s*(.*)", rest)
                    if lm:
                        label = lm.group(1) or None
                        rest = lm.group(2).strip()
                if not src or not rest:
                    continue
                nodes.setdefault(src, src.replace("_", " "))
                nodes.setdefault(rest, rest.replace("_", " "))
                if src not in order:
                    order.append(src)
                if rest not in order:
                    order.append(rest)
                edges.append((src, rest, label))
        return nodes, order, edges

    def _layout_diagram(self, nodes, order, edges, layout):
        ids = [i for i in order if i in nodes] or list(nodes)
        indeg = {i: 0 for i in ids}
        adj = {i: [] for i in ids}
        for a, b, _ in edges:
            if a in adj and b in indeg:
                adj[a].append(b)
                indeg[b] += 1
        q = deque([i for i in ids if indeg.get(i, 0) == 0])
        if not q:
            q = deque([ids[0]])
        level = {i: 0 for i in ids}
        seen = set()
        while q:
            a = q.popleft()
            if a in seen:
                continue
            seen.add(a)
            for b in adj[a]:
                level[b] = max(level[b], level[a] + 1)
                indeg[b] -= 1
                if indeg[b] <= 0:
                    q.append(b)
        groups = {}
        for i in ids:
            groups.setdefault(level.get(i, 0), []).append(i)
        return groups

    def _diagram_positions(self, groups, layout):
        positions = {}
        levels = sorted(groups)
        n = len(levels)
        bx, by, bw, bh = 1.0, 1.95, 11.3, 4.9
        if layout.upper() in ("LR", "RL"):
            col_w = bw / n
            for li, lvl in enumerate(levels):
                members = groups[lvl]
                cx = bx + col_w * li + col_w / 2
                step = bh / max(len(members), 1)
                for ni, nid in enumerate(members):
                    cy = by + step * ni + step / 2
                    positions[nid] = (cx, cy)
        else:
            row_h = bh / n
            for li, lvl in enumerate(levels):
                members = groups[lvl]
                cy = by + row_h * li + row_h / 2
                step = bw / max(len(members), 1)
                for ni, nid in enumerate(members):
                    cx = bx + step * ni + step / 2
                    positions[nid] = (cx, cy)
        return positions

    def _content_header(self, heading):
        slide = self._blank()
        self._bg(slide, COLORS["paper"])
        self._rect(slide, 0, 0, SLIDE_W, 0.07, COLORS["primary"], MSO_SHAPE.RECTANGLE)
        self._text(slide, 0.72, 0.34, 11.8, 0.72, heading,
                   size=24, color=COLORS["dark"], bold=True, align=PP_ALIGN.LEFT)
        self._rect(slide, 0.74, 1.2, 1.5, 0.045, COLORS["accent"], MSO_SHAPE.RECTANGLE)
        self._rect(slide, 0, 7.26, SLIDE_W, 0.05, COLORS["light"], MSO_SHAPE.RECTANGLE)
        self._footer(slide, "PaperWise")
        return slide

    def _bullets(self, slide, items, x, y, w, h, size=18):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        n = len(items)
        fs = size if n <= 4 else (size - 2 if n <= 6 else size - 4)
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "▪  " + str(item)
            p.space_after = Pt(12)
            self._run_font(p.runs[0] if p.runs else p.add_run(), fs, COLORS["dark"], False)

    def _stat_cards(self, slide, stats, x, y, w, h):
        if not stats:
            return
        per = h / max(len(stats), 1)
        for i, st in enumerate(stats[:4]):
            card = self._rect(slide, x, y + i * per, w, per - 0.25, COLORS["light"])
            tf = card.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.2)
            tf.margin_top = Inches(0.12)
            p = tf.paragraphs[0]
            p.text = str(st.get("value", ""))
            self._run_font(p.runs[0] if p.runs else p.add_run(), 30, COLORS["primary2"], True)
            p2 = tf.add_paragraph()
            p2.text = str(st.get("label", ""))
            self._run_font(p2.runs[0] if p2.runs else p2.add_run(), 14, COLORS["gray"], False)

    def _key_point_badge(self, slide, text, x, y, w, h):
        card = self._rect(slide, x, y, w, h, COLORS["accent_soft"])
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)
        p = tf.paragraphs[0]
        p.text = "关键结论"
        self._run_font(p.runs[0] if p.runs else p.add_run(), 14, COLORS["accent"], True)
        p2 = tf.add_paragraph()
        p2.text = text
        p2.space_before = Pt(8)
        self._run_font(p2.runs[0] if p2.runs else p2.add_run(), 18, COLORS["dark"], True)

    def _card_column(self, slide, x, y, w, h, title, items, color, soft):
        card = self._rect(slide, x, y, w, h, soft)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.18)
        p = tf.paragraphs[0]
        p.text = title
        self._run_font(p.runs[0] if p.runs else p.add_run(), 20, color, True)
        for item in items[:5]:
            pp = tf.add_paragraph()
            pp.text = "• " + str(item)
            pp.space_after = Pt(8)
            self._run_font(pp.runs[0] if pp.runs else pp.add_run(), 14, COLORS["dark"], False)

    def _picture_fit(self, slide, path):
        if Image:
            try:
                with Image.open(str(path)) as im:
                    iw, ih = im.size
            except Exception:
                iw = ih = 1
        else:
            iw = ih = 1
        iw, ih = iw or 1, ih or 1
        top, bottom = Inches(1.6), Inches(7.1)
        max_w, max_h = Inches(11.3), bottom - top
        ratio = min(int(max_w) / iw, int(max_h) / ih)
        w, h = int(iw * ratio), int(ih * ratio)
        left = int((self.prs.slide_width - w) / 2)
        top_px = int(top + max(0, (max_h - h) / 2))
        slide.shapes.add_picture(str(path), left, top_px, width=w, height=h)

    def _footer(self, slide, text):
        self._text(slide, 0.5, 7.08, 9.0, 0.35, text, size=9, color=COLORS["gray"])
        self._text(slide, 11.9, 7.08, 0.9, 0.35, str(self._slide_no),
                   size=9, color=COLORS["gray"], align=PP_ALIGN.RIGHT)

    def _text(self, slide, x, y, w, h, text, size=18, color=COLORS["dark"],
              bold=False, align=PP_ALIGN.LEFT, italic=False):
        box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        lines = str(text).split("\n")
        for i, line in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.alignment = align
            self._run_font(p.runs[0] if p.runs else p.add_run(), size, color, bold, italic)
        return box

    def _style_cell(self, cell, bold, color, size):
        tf = cell.text_frame
        tf.word_wrap = True
        for p in tf.paragraphs:
            for run in p.runs:
                self._run_font(run, size, color, bold)

    @staticmethod
    def _run_font(run, size, color, bold=False, italic=False):
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = FONT_BODY
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn("a:ea"))
        if ea is None:
            ea = rPr.makeelement(qn("a:ea"), {})
            rPr.append(ea)
        ea.set("typeface", FONT_BODY)


def render_deck(deck: dict, output_path: str) -> str:
    """便捷入口：渲染 deck JSON 到 .pptx。"""
    return SlideDeckRenderer().render(deck, output_path)
