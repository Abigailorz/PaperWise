"""PPT 生成器 — 基于 python-pptx 生成学术汇报演示文稿（v2 重写）

v2 相比旧版修复的问题：
- 章节内容正确映射，空章节自动跳过，不再生成 "(No content available)" 空页
- 自动嵌入论文解析出的图片（figures/）与表格（tables/）
- 更稳健的要点提取：去除 Markdown/引用噪声、按句切分、控制长度、去重
- 文本自适应字号，避免溢出
- 统一中英文字体与更现代的学术配色
"""

from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

try:
    from PIL import Image
except Exception:  # pragma: no cover
    Image = None


class PPTXGenerator:
    """学术论文演示文稿生成器。"""

    COLORS = {
        "primary": RGBColor(0x0F, 0x2A, 0x5C),
        "primary_2": RGBColor(0x1E, 0x4E, 0xB8),
        "accent": RGBColor(0x3B, 0x82, 0xF6),
        "accent_soft": RGBColor(0xDB, 0xE8, 0xFE),
        "orange": RGBColor(0xF5, 0x9E, 0x0B),
        "dark": RGBColor(0x1F, 0x29, 0x37),
        "gray": RGBColor(0x6B, 0x72, 0x80),
        "light": RGBColor(0xF3, 0xF6, 0xFB),
        "white": RGBColor(0xFF, 0xFF, 0xFF),
        "green": RGBColor(0x10, 0x98, 0x5A),
        "red": RGBColor(0xD9, 0x43, 0x52),
    }

    FONT_TITLE = "Microsoft YaHei"
    FONT_BODY = "Microsoft YaHei"

    SECTION_ORDER = [
        ("overview", "论文概览"),
        ("motivation", "研究动机与问题"),
        ("methodology", "核心方法"),
        ("experiments", "实验与结果"),
        ("critical_analysis", "批判性分析"),
        ("related_work", "相关工作"),
        ("conclusion", "总结与展望"),
    ]

    def __init__(self, workspace: Path):
        self.workspace = Path(workspace)
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9
        self.prs.slide_height = Inches(7.5)

    # ============ 主入口 ============

    def generate(self, paper_data: dict, output_path: Optional[str] = None) -> str:
        sections = dict(paper_data.get("sections") or {})
        for key, _ in self.SECTION_ORDER:
            if key in paper_data and key not in sections:
                sections[key] = paper_data[key]

        title = paper_data.get("title") or self.workspace.name
        authors = paper_data.get("authors") or paper_data.get("author") or ""
        venue = paper_data.get("venue") or paper_data.get("subject") or ""
        year = paper_data.get("year") or ""

        figures = list(paper_data.get("figures") or self._discover_figures())
        tables = list(paper_data.get("tables") or self._discover_tables())

        self._add_title_slide(title, authors, venue, year)

        for key, section_title in self.SECTION_ORDER:
            content = (sections.get(key) or "").strip()
            if not content:
                continue
            if key == "critical_analysis":
                self._add_critical_slide(content)
            else:
                self._add_section_slides(section_title, content)

        for fig in figures[:4]:
            self._add_figure_slide(fig)

        if tables:
            self._add_table_slide(tables[0])

        self._add_thank_you_slide()

        output = output_path or str(self.workspace / "presentation" / "slides.pptx")
        Path(output).parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(output)
        return output

    # ============ 幻灯片构建 ============

    def _add_title_slide(self, title, authors, venue, year):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide, self.COLORS["primary"])

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), self.prs.slide_width, Inches(0.12)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.COLORS["orange"]
        bar.line.fill.background()

        box = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(2.6))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(38)
        p.font.bold = True
        p.font.color.rgb = self.COLORS["white"]
        p.font.name = self.FONT_TITLE
        p.alignment = PP_ALIGN.LEFT

        sub = slide.shapes.add_textbox(Inches(0.9), Inches(4.6), Inches(11.5), Inches(1.4))
        tf = sub.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = authors
        p.font.size = Pt(19)
        p.font.color.rgb = RGBColor(0xBF, 0xD1, 0xF5)
        p.font.name = self.FONT_BODY
        p2 = tf.add_paragraph()
        p2.text = " · ".join(x for x in [venue, year] if x)
        p2.font.size = Pt(15)
        p2.font.color.rgb = RGBColor(0x9A, 0xB0, 0xD8)
        p2.font.name = self.FONT_BODY
        p2.space_before = Pt(10)

        self._add_footer(slide, "PaperWise · AI 论文解读")

    def _add_section_slides(self, title, content):
        points = self._extract_points(content, max_points=12)
        if not points:
            return
        per_page = 5
        for i in range(0, len(points), per_page):
            chunk = points[i:i + per_page]
            page_title = title if i == 0 else f"{title}（续）"
            slide = self._create_content_slide(page_title)
            self._add_bullet_list(slide, chunk, Inches(0.9), Inches(1.7), Inches(11.5), Inches(5.2))

    def _add_critical_slide(self, content):
        strengths, limitations = self._split_strengths_limitations(content)
        slide = self._create_content_slide("批判性分析")
        self._add_two_column(slide, Inches(0.8), "✓ 优势", strengths, self.COLORS["green"])
        self._add_two_column(slide, Inches(6.9), "▲ 局限", limitations, self.COLORS["red"])

    def _add_figure_slide(self, fig):
        path = fig.get("path") or fig.get("file") or ""
        if not path:
            return
        p = Path(path)
        if not p.is_absolute():
            p = self.workspace / p
        if not p.exists():
            return
        slide = self._create_content_slide("论文图示")
        self._add_picture_fit(slide, p)
        caption = fig.get("caption") or ""
        if caption:
            cb = slide.shapes.add_textbox(Inches(1.0), Inches(6.7), Inches(11.3), Inches(0.5))
            tf = cb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = caption[:160]
            p.font.size = Pt(12)
            p.font.italic = True
            p.font.color.rgb = self.COLORS["gray"]
            p.font.name = self.FONT_BODY
            p.alignment = PP_ALIGN.CENTER

    def _add_table_slide(self, table):
        headers = table.get("headers") or []
        rows = table.get("rows") or []
        if not headers:
            return
        n_cols = min(len(headers), 6)
        n_rows = min(len(rows) + 1, 8)
        slide = self._create_content_slide("关键表格")
        frame = slide.shapes.add_table(
            n_rows, n_cols, Inches(0.7), Inches(1.7), Inches(12.0), Inches(0.5 * n_rows)
        )
        gtbl = frame.table
        for c, h in enumerate(headers[:n_cols]):
            cell = gtbl.cell(0, c)
            cell.text = str(h)[:60]
            cell.fill.solid()
            cell.fill.fore_color.rgb = self.COLORS["primary"]
            self._style_cell(cell, bold=True, color=self.COLORS["white"], size=13)
        for r, row in enumerate(rows[:n_rows - 1], start=1):
            for c in range(n_cols):
                val = row[c] if c < len(row) else ""
                cell = gtbl.cell(r, c)
                cell.text = str(val)[:80]
                cell.fill.solid()
                cell.fill.fore_color.rgb = self.COLORS["light"] if r % 2 else self.COLORS["white"]
                self._style_cell(cell, bold=False, color=self.COLORS["dark"], size=12)

    def _add_thank_you_slide(self):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide, self.COLORS["primary"])
        box = slide.shapes.add_textbox(Inches(1), Inches(2.7), Inches(11.3), Inches(2))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = "谢谢观看"
        p.font.size = Pt(48)
        p.font.bold = True
        p.font.color.rgb = self.COLORS["white"]
        p.font.name = self.FONT_TITLE
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = "Q & A"
        p2.font.size = Pt(24)
        p2.font.color.rgb = RGBColor(0xBF, 0xD1, 0xF5)
        p2.font.name = self.FONT_BODY
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(18)
        self._add_footer(slide, "由 PaperWise 生成")

    # ============ 辅助 ============

    def _create_content_slide(self, title):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self._set_bg(slide, self.COLORS["white"])

        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), self.prs.slide_width, Inches(1.15)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = self.COLORS["primary"]
        bar.line.fill.background()

        tf = bar.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.55)
        tf.margin_top = Inches(0.18)
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = self.COLORS["white"]
        p.font.name = self.FONT_TITLE
        p.alignment = PP_ALIGN.LEFT

        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(7.2), self.prs.slide_width, Inches(0.06)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = self.COLORS["orange"]
        line.line.fill.background()

        self._add_footer(slide, self._slide_number(slide))
        return slide

    def _add_bullet_list(self, slide, items, left, top, width, height):
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.TOP
        size = self._fit_size(items)
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(size)
            p.font.color.rgb = self.COLORS["dark"]
            p.font.name = self.FONT_BODY
            p.space_after = Pt(14)
            p.level = 0

    def _add_two_column(self, slide, left, heading, items, color):
        box = slide.shapes.add_textbox(left, Inches(1.7), Inches(5.6), Inches(5.0))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = heading
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = self.FONT_TITLE
        p.space_after = Pt(10)
        for item in items[:4]:
            p = tf.add_paragraph()
            p.text = "• " + item
            p.font.size = Pt(14)
            p.font.color.rgb = self.COLORS["dark"]
            p.font.name = self.FONT_BODY
            p.space_after = Pt(8)

    def _add_picture_fit(self, slide, path):
        if Image:
            try:
                with Image.open(str(path)) as im:
                    iw, ih = im.size
            except Exception:
                iw = ih = 1
        else:
            iw = ih = 1
        iw = iw or 1
        ih = ih or 1
        top_in = Inches(1.6)
        bottom_in = Inches(7.15)
        max_w = Inches(11.3)
        max_h = bottom_in - top_in
        ratio = min(int(max_w) / iw, int(max_h) / ih)
        w = int(iw * ratio)
        h = int(ih * ratio)
        left = int((self.prs.slide_width - w) / 2)
        top = int(top_in + max(0, (max_h - h) / 2))
        slide.shapes.add_picture(str(path), left, top, width=w, height=h)

    def _add_footer(self, slide, text):
        footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.05), Inches(9), Inches(0.4))
        tf = footer.text_frame
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(9)
        p.font.color.rgb = self.COLORS["gray"]
        p.font.name = self.FONT_BODY

    def _set_bg(self, slide, color):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def _slide_number(self, slide):
        try:
            return f"{list(self.prs.slides).index(slide) + 1}"
        except ValueError:
            return ""

    def _style_cell(self, cell, bold, color, size):
        tf = cell.text_frame
        tf.word_wrap = True
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.bold = bold
                run.font.size = Pt(size)
                run.font.color.rgb = color
                run.font.name = self.FONT_BODY

    # ============ 内容提取 ============

    def _extract_points(self, text, max_points=8):
        import re
        if not text:
            return []
        text = re.sub(r"```.*?```", " ", text, flags=re.DOTALL)
        text = re.sub(r"!\[[^\]]*\]\([^)]*\)", " ", text)
        text = re.sub(r"\[([^\]]+)\]\([^)]*\)", r"\1", text)
        text = re.sub(r"#{1,6}\s*", "", text)
        text = text.replace("\r", "\n")

        lines = [l.strip() for l in text.split("\n") if l.strip()]
        raw = []
        for line in lines:
            low = line.lower()
            if low.startswith(("figure ", "table ", "fig.", "tab.")):
                continue
            if len(line) > 260:
                raw.extend(re.split(r"(?<=[。.!?;；])\s*", line))
            else:
                raw.append(line)

        points = []
        seen = set()
        for r in raw:
            r = re.sub(r"^\s*(?:[-*•·>]|\d+[.、)])\s*", "", r).strip()
            r = re.sub(r"\s+", " ", r)
            if not (8 <= len(r) <= 220):
                continue
            key = r[:60]
            if key in seen:
                continue
            seen.add(key)
            if len(r) > 220:
                r = r[:219] + "…"
            points.append(r)
            if len(points) >= max_points * 2:
                break
        return points[:max_points]

    def _split_strengths_limitations(self, content):
        import re
        low = content.lower()
        markers = [
            ("limitation", ["limitations", "limitation", "局限", "不足", "缺点", "weakness"]),
            ("strength", ["strengths", "strength", "优势", "亮点", "优点", "contributions"]),
        ]
        split_idx = None
        split_kind = None
        for kind, keys in markers:
            for k in keys:
                i = low.find(k)
                if i >= 0 and (split_idx is None or i < split_idx):
                    split_idx = i
                    split_kind = kind

        strengths, limitations = [], []
        if split_idx is not None:
            part_a = content[:split_idx]
            part_b = content[split_idx:]
            if split_kind == "limitation":
                strengths, limitations = part_a, part_b
            else:
                strengths, limitations = part_b, part_a
        else:
            limitations = content

        def to_points(s):
            pts = self._extract_points(s, max_points=4)
            return pts if pts else (["（内容待补充）"] if s.strip() else [])

        return to_points(strengths), to_points(limitations)

    def _discover_figures(self):
        figs = []
        fig_dir = self.workspace / "figures"
        if fig_dir.exists():
            for ext in ("*.png", "*.jpg", "*.jpeg"):
                for f in sorted(fig_dir.glob(ext)):
                    figs.append({"path": str(f), "caption": ""})
        return figs

    def _discover_tables(self):
        import json
        tables = []
        tbl_dir = self.workspace / "tables"
        if tbl_dir.exists():
            for f in sorted(tbl_dir.glob("table_*.json")):
                try:
                    tables.append(json.loads(f.read_text(encoding="utf-8")))
                except Exception:
                    continue
        return tables

    def _fit_size(self, items):
        max_len = max((len(i) for i in items), default=0)
        n = len(items)
        if n <= 3 and max_len < 60:
            return 20
        if n <= 4 and max_len < 120:
            return 18
        if n <= 5 and max_len < 160:
            return 16
        return 14


def get_pptx_system_prompt() -> str:
    """PPT 生成的 Agent 系统提示词。"""
    return """You are a presentation designer specialized in academic paper presentations.

<your_task>
Create a professional PowerPoint presentation (10-15 slides) summarizing an academic paper.
The presentation should be clear, well-structured, and suitable for a research seminar.
</your_task>

<slide_structure>
1. Title Slide (1 slide)
2. Overview (1 slide) — One-paragraph summary + key contributions
3. Background & Motivation (1-2 slides) — Problem context, prior work limitations
4. Problem Statement (1 slide) — Clear research question
5. Core Method (2-3 slides) — Key ideas, architecture, important equations
6. Experimental Results (2-3 slides) — Setup, main results, key comparisons
7. Critical Analysis (1 slide) — Strengths vs Limitations
8. Conclusion (1 slide) — Summary + future work
9. Thank You / Q&A (1 slide)
</slide_structure>

<design_guidelines>
- Each slide should have a clear single message
- Use bullet points, not paragraphs
- Include key numbers and metrics
- Reference specific figures/tables from the paper
- Keep text concise — slides support the talk, not replace it
</design_guidelines>
"""
